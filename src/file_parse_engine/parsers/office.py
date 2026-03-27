"""Office document parsers (DOCX / PPTX) — convert to PDF via LibreOffice, then render pages."""

from __future__ import annotations

import io
import shutil
import subprocess
import tempfile
from pathlib import Path

import fitz  # PyMuPDF

from file_parse_engine.config import get_settings
from file_parse_engine.models import PageImage, ParsedDocument, ParsedPage
from file_parse_engine.parsers import register
from file_parse_engine.parsers.base import BaseParser
from file_parse_engine.utils.logger import get_logger

logger = get_logger("parsers.office")


def _find_libreoffice() -> str | None:
    """Locate the LibreOffice binary."""
    for name in ("libreoffice", "soffice"):
        path = shutil.which(name)
        if path:
            return path

    # macOS common paths
    mac_paths = [
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
    ]
    for p in mac_paths:
        if Path(p).exists():
            return p

    return None


def _convert_to_pdf(source: Path) -> bytes:
    """Convert an Office document to PDF using LibreOffice headless."""
    lo = _find_libreoffice()
    if lo is None:
        raise RuntimeError(
            "LibreOffice not found. Install it for DOCX/PPTX support:\n"
            "  macOS: brew install --cask libreoffice\n"
            "  Linux: sudo apt install libreoffice-common"
        )

    with tempfile.TemporaryDirectory() as tmpdir:
        result = subprocess.run(
            [lo, "--headless", "--convert-to", "pdf", "--outdir", tmpdir, str(source)],
            capture_output=True,
            text=True,
            timeout=120,
        )

        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

        # Find the output PDF
        pdf_files = list(Path(tmpdir).glob("*.pdf"))
        if not pdf_files:
            raise RuntimeError("LibreOffice produced no PDF output")

        return pdf_files[0].read_bytes()


def _pdf_bytes_to_page_images(pdf_bytes: bytes, dpi: int) -> list[PageImage]:
    """Render PDF bytes to page images."""
    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    pages: list[PageImage] = []

    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(matrix=matrix)

        buf = io.BytesIO()
        buf.write(pix.tobytes("png"))
        image_bytes = buf.getvalue()

        pages.append(PageImage(
            page_number=page_num + 1,
            image_bytes=image_bytes,
            width=pix.width,
            height=pix.height,
        ))

    doc.close()
    return pages


# ---------------------------------------------------------------------------
# DOCX Parser
# ---------------------------------------------------------------------------


@register("docx")
class DOCXParser(BaseParser):
    """DOCX parser: LibreOffice → PDF → page images → VLM."""

    file_type = "docx"

    async def to_page_images(self, path: Path) -> list[PageImage]:
        settings = get_settings()
        logger.info("Converting DOCX to PDF: %s", path.name)

        try:
            pdf_bytes = _convert_to_pdf(path)
        except RuntimeError:
            logger.warning("LibreOffice not available, falling back to programmatic extraction")
            return await self._fallback_extraction(path)

        pages = _pdf_bytes_to_page_images(pdf_bytes, settings.image_dpi)
        logger.info("DOCX rendered: %d page(s)", len(pages))
        return pages

    async def _fallback_extraction(self, path: Path) -> list[PageImage]:
        """Fallback: extract text via python-docx and render to image.

        If LibreOffice is unavailable, we extract text directly and
        let the engine handle it via the direct path instead.
        """
        raise RuntimeError(
            "DOCX VLM extraction requires LibreOffice. "
            "Install LibreOffice or use programmatic extraction."
        )

    @property
    def needs_vlm(self) -> bool:
        return _find_libreoffice() is not None

    async def to_markdown_direct(self, path: Path) -> ParsedDocument:
        """Fallback: extract text programmatically via python-docx."""
        from docx import Document

        doc = Document(str(path))
        parts: list[str] = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = (para.style.name or "").lower()
            if "heading 1" in style_name:
                parts.append(f"# {text}")
            elif "heading 2" in style_name:
                parts.append(f"## {text}")
            elif "heading 3" in style_name:
                parts.append(f"### {text}")
            elif "heading" in style_name:
                parts.append(f"#### {text}")
            elif "list" in style_name:
                parts.append(f"- {text}")
            else:
                parts.append(text)

        # Extract tables
        for table in doc.tables:
            parts.append("")
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)

            if rows:
                # Header
                parts.append("| " + " | ".join(rows[0]) + " |")
                parts.append("| " + " | ".join("---" for _ in rows[0]) + " |")
                for row in rows[1:]:
                    parts.append("| " + " | ".join(row) + " |")
            parts.append("")

        return ParsedDocument(
            source_path=str(path),
            file_type="docx",
            pages=[ParsedPage(page_number=1, markdown="\n".join(parts))],
        )

    def extract_metadata(self, path: Path) -> dict:
        from docx import Document

        doc = Document(str(path))
        props = doc.core_properties
        return {
            "title": props.title or "",
            "author": props.author or "",
            "subject": props.subject or "",
            "created": str(props.created) if props.created else "",
            "modified": str(props.modified) if props.modified else "",
        }


# ---------------------------------------------------------------------------
# PPTX Parser
# ---------------------------------------------------------------------------


@register("pptx")
class PPTXParser(BaseParser):
    """PPTX parser: LibreOffice → PDF → page images → VLM."""

    file_type = "pptx"

    async def to_page_images(self, path: Path) -> list[PageImage]:
        settings = get_settings()
        logger.info("Converting PPTX to PDF: %s", path.name)

        try:
            pdf_bytes = _convert_to_pdf(path)
        except RuntimeError:
            logger.warning("LibreOffice not available, falling back to programmatic extraction")
            raise

        pages = _pdf_bytes_to_page_images(pdf_bytes, settings.image_dpi)
        logger.info("PPTX rendered: %d slide(s)", len(pages))
        return pages

    @property
    def needs_vlm(self) -> bool:
        return _find_libreoffice() is not None

    async def to_markdown_direct(self, path: Path) -> ParsedDocument:
        """Fallback: extract text programmatically via python-pptx."""
        from pptx import Presentation

        prs = Presentation(str(path))
        pages: list[ParsedPage] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            parts: list[str] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            # First text frame is usually the title
                            if shape == slide.shapes[0] and not parts:
                                parts.append(f"## {text}")
                            else:
                                parts.append(f"- {text}" if paragraph.level > 0 else text)

                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)

                    if rows:
                        parts.append("")
                        parts.append("| " + " | ".join(rows[0]) + " |")
                        parts.append("| " + " | ".join("---" for _ in rows[0]) + " |")
                        for row in rows[1:]:
                            parts.append("| " + " | ".join(row) + " |")
                        parts.append("")

            pages.append(ParsedPage(
                page_number=slide_num,
                markdown="\n".join(parts),
            ))

        return ParsedDocument(
            source_path=str(path),
            file_type="pptx",
            pages=pages,
        )

    def extract_metadata(self, path: Path) -> dict:
        from pptx import Presentation

        prs = Presentation(str(path))
        props = prs.core_properties
        return {
            "title": props.title or "",
            "author": props.author or "",
            "subject": props.subject or "",
            "slide_count": len(prs.slides),
        }
